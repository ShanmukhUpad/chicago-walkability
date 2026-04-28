"""
Build Undergraduate_Research_Symposium.pptx
Run after knitting finalProject.Rmd so poster_figures/ is populated.
Requires: pip install python-pptx Pillow
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
NAVY       = RGBColor(0x00, 0x33, 0x66)
ORANGE     = RGBColor(0xFF, 0x5F, 0x05)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT  = RGBColor(0x22, 0x22, 0x22)
MID_TEXT   = RGBColor(0x55, 0x55, 0x55)

# ---------------------------------------------------------------------------
# Poster geometry (inches)
# ---------------------------------------------------------------------------
PW, PH = 48.0, 36.0

HEADER_H = 2.6
BODY_Y   = HEADER_H + 0.2
BODY_H   = PH - BODY_Y - 0.35

MARGIN  = 0.4
GAP     = 0.35
C1W     = 11.0
C2W     = 13.5
C1X     = MARGIN
C2X     = C1X + C1W + GAP
C3X     = C2X + C2W + GAP
C3W     = PW - MARGIN - C3X   # ~22.4"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def I(n):
    return Inches(n)

def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def section_header(slide, x, y, w, text, fs=24):
    bar = rect(slide, x, y, w, 0.42, fill=NAVY)
    tf = bar.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r  = p.add_run(); r.text = "  " + text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE
    return bar

def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tb.word_wrap = True
    tb.text_frame.word_wrap = True
    return tb.text_frame

def para(tf, text, fs=18, bold=False, color=DARK_TEXT,
         align=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(0),
         italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after  = space_after
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def first_para(tf, text, fs=18, bold=False, color=DARK_TEXT,
               align=PP_ALIGN.LEFT, space_before=Pt(0)):
    p = tf.paragraphs[0]; p.alignment = align; p.space_before = space_before
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = color
    return p

def image(slide, path, x, y, w=None, h=None):
    """Insert image. Specify w XOR h to preserve aspect ratio."""
    path = BASE / path if not Path(path).is_absolute() else Path(path)
    if not path.exists():
        b = rect(slide, x, y, w or 4, h or 3, fill=LIGHT_GRAY, line=MID_TEXT)
        tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"[{path.name}]"
        r.font.size = Pt(14); r.font.color.rgb = MID_TEXT
        return
    kwargs = {}
    if w is not None: kwargs['width']  = I(w)
    if h is not None: kwargs['height'] = I(h)
    slide.shapes.add_picture(str(path), I(x), I(y), **kwargs)

def caption(slide, x, y, w, text, fs=14):
    tf = textbox(slide, x, y, w, 0.35)
    first_para(tf, text, fs=fs, color=MID_TEXT, align=PP_ALIGN.CENTER, bold=True)

def stat_box(slide, x, y, w, h, number, label, ns=34, ls=16):
    b = rect(slide, x, y, w, h, fill=WHITE, line=ORANGE, lw=Pt(2.5))
    tf = b.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r  = p.add_run(); r.text = number
    r.font.size = Pt(ns); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(ls); r2.font.color.rgb = MID_TEXT


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = I(PW)
prs.slide_height = I(PH)
slide = prs.slides.add_slide(prs.slide_layouts[6])

rect(slide, 0, 0, PW, PH, fill=WHITE)

# ── HEADER ──────────────────────────────────────────────────────────────────
rect(slide, 0, 0, PW, 0.45, fill=NAVY)
rect(slide, 0, 0.45, PW, 0.14, fill=ORANGE)

title_tf = textbox(slide, 0.5, 0.65, 45, 1.35)
first_para(title_tf,
    "A Spatial Analysis of Walkability and Urban Form in the Chicago Metropolitan Area",
    fs=40, bold=True, color=NAVY)

auth_tf = textbox(slide, 0.5, 2.05, 45, 0.55)
first_para(auth_tf,
    "Shanmukh Upadhyayula  |  Department of Geography & GIS, University of Illinois Urbana-Champaign  |  supad83@illinois.edu",
    fs=22, color=MID_TEXT)

rect(slide, 0, HEADER_H, PW, 0.1, fill=ORANGE)


# ── COLUMN 1: Introduction, Hypothesis, Methods, Score Table, References ────
x1, y1 = C1X, BODY_Y

section_header(slide, x1, y1, C1W, "Introduction")
y1 += 0.46
tf = textbox(slide, x1, y1, C1W, 6.5)
first_para(tf,
    "Walkability measures how safe an area is for pedestrians and how easily amenities "
    "can be reached on foot. The ChiVes database, which aggregates Chicago environmental "
    "data for public use, lacked a walkability measure for the built environment. This "
    "project integrates the EPA National Walkability Index (2021), originally at the "
    "block group scale and converted to 2020 census tracts for ChiVes compatibility, "
    "to visualize how walkability varies across Chicago and to assess its relationship "
    "with transit access, employment diversity, and economic opportunity. Chicago ranks "
    "high nationally but outer, car-centric neighborhoods remain underserved. Improving "
    "walkability has direct implications for urban sustainability and transportation equity.",
    fs=17, color=DARK_TEXT)
y1 += 6.7

section_header(slide, x1, y1, C1W, "Hypothesis")
y1 += 0.46
tf = textbox(slide, x1, y1, C1W, 3.5)
first_para(tf,
    "Walkability is expected to be highest in the Loop and CBD and to decline toward "
    "outer neighborhoods. Because Chicago is the third most populous U.S. city, most "
    "census tracts should fall in the above-average tier, with a meaningful share "
    "reaching the highly walkable range. Population density and transit access are "
    "included as exploratory variables to assess what actually drives the walkability "
    "score at the tract level.",
    fs=17, color=DARK_TEXT)
y1 += 3.7

section_header(slide, x1, y1, C1W, "Methods")
y1 += 0.46
tf = textbox(slide, x1, y1, C1W, 7.5)
first_para(tf,
    "Summary statistics (mean, median, standard deviation, and range of walkability "
    "scores, alongside mean population density and transit distance) provide a baseline "
    "description of the data. Spatial autocorrelation is tested using Moran's I to "
    "determine whether high- and low-walkability tracts cluster together rather than "
    "occurring randomly.\n\n"
    "A Pearson correlation matrix and linear regression quantify the relationship between "
    "walkability and each of the four EPA inputs: employment + household entropy (D2A), "
    "8-tier employment mix (D2B), street intersection density (D3B), and distance to the "
    "nearest transit stop (D4A). The regression is fit on City of Chicago tracts "
    "(R² = 0.85). Analysis of Variance (ANOVA) compares mean EPA-input values across "
    "the five walkability categories. Local Moran's I (LISA) identifies statistically "
    "significant spatial clusters of high- and low-walkability tracts (p < 0.05).",
    fs=17, color=DARK_TEXT)
y1 += 7.7

# Walkability score categories table (color-coded)
section_header(slide, x1, y1, C1W, "Walkability Score Categories", fs=22)
y1 += 0.46
cat_rows = [
    ("Category",           "Score Range",    "Metro Tracts",    NAVY,            WHITE),
    ("Least Walkable",     "1.00 – 5.75",    "53 (2.6%)",       RGBColor(0x44,0x01,0x54), WHITE),
    ("Below Average",      "5.76 – 10.50",   "496 (24.6%)",     RGBColor(0x31,0x68,0x8E), WHITE),
    ("Above Average",      "10.51 – 15.25",  "1,090 (54.1%)",   RGBColor(0x35,0xB7,0x79), DARK_TEXT),
    ("Highly Walkable",    "15.26 – 18.09",  "350 (17.4%)",     RGBColor(0x90,0xD7,0x43), DARK_TEXT),
    ("Most Walkable",      "18.10 – 20.00",  "20 (1.0%)",       RGBColor(0xFD,0xE7,0x25), DARK_TEXT),
]
row_h = 0.54
col_ws = [4.2, 3.3, 3.2]
col_xs = [x1, x1 + col_ws[0], x1 + col_ws[0] + col_ws[1]]
for ri, (cat, rng, cnt, bg, fg) in enumerate(cat_rows):
    ry = y1 + ri * row_h
    is_hdr = ri == 0
    for ci, (txt, cw_col, col_x) in enumerate(zip([cat, rng, cnt], col_ws, col_xs)):
        cell = rect(slide, col_x, ry, cw_col, row_h,
                    fill=bg if is_hdr else bg,
                    line=NAVY, lw=Pt(0.75))
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ("  " if ci == 0 else "") + txt
        r.font.size = Pt(16 if not is_hdr else 17)
        r.font.bold = is_hdr
        r.font.color.rgb = fg if not is_hdr else WHITE
y1 += len(cat_rows) * row_h + 0.3

# References
section_header(slide, x1, y1, C1W, "References", fs=21)
y1 += 0.46
tf = textbox(slide, x1, y1, C1W, 1.6)
first_para(tf,
    "1. U.S. EPA (2021). National Walkability Index: User Guide and Methodology.\n"
    "2. U.S. EPA (2021). Smart Location Database Technical Documentation.\n"
    "3. ChiVes Project Team (2024). ChiVes Platform. UIUC.",
    fs=14, color=MID_TEXT)
y1 += 1.7

# Acknowledgements
section_header(slide, x1, y1, C1W, "Acknowledgements", fs=21)
y1 += 0.46
tf = textbox(slide, x1, y1, C1W, 0.8)
first_para(tf,
    "Dr. Marynia Kolak, Department of Geography & GIS, University of Illinois Urbana-Champaign",
    fs=14, color=MID_TEXT)


# ── COLUMN 2: How Walkability is Calculated + Key Stats + Corrplot ──────────
x2, y2 = C2X, BODY_Y

section_header(slide, x2, y2, C2W, "How Walkability is Calculated (EPA NatWalkInd)")
y2 += 0.46

calc_h = 7.2
calc_box = rect(slide, x2, y2, C2W, calc_h, fill=LIGHT_BLUE, line=NAVY, lw=Pt(1.5))
tf = calc_box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  Formula:  NatWalkInd = ¼ × (D2Aₚ + D2Bₚ + D3Bₚ + D4Aₚ)"
r.font.size = Pt(21); r.font.bold = True; r.font.color.rgb = NAVY

p2 = tf.add_paragraph(); p2.space_before = Pt(4)
r2 = p2.add_run()
r2.text = "  Each component is ranked 1–20 nationally, then averaged. Final score: 1 (least) → 20 (most walkable)."
r2.font.size = Pt(17); r2.font.color.rgb = DARK_TEXT

# Component table inside the calc box
table_rows = [
    ("Component",   "Description",                                      "City r", True),
    ("D2A_EPHHM",  "Employment + Household Entropy (mixed-use intensity)", "0.69",  False),
    ("D2B_E8MIXA", "8-Tier Employment Mix (job diversity)",               "0.68",  False),
    ("D3B",        "Street Intersection Density (network connectivity)",  "0.53",  False),
    ("D4A",        "Distance to Nearest Transit Stop",                    "−0.25", False),
]
tr_h  = 0.82
tr_y0 = y2 + 1.55
pad   = 0.15
cws   = [2.8, 8.0, 1.8]
cxs   = [x2 + pad, x2 + pad + cws[0], x2 + pad + cws[0] + cws[1]]
for ri, (comp, desc, corr, hdr) in enumerate(table_rows):
    ry  = tr_y0 + ri * tr_h
    bg  = NAVY if hdr else (LIGHT_BLUE if ri % 2 == 0 else WHITE)
    fg  = WHITE if hdr else DARK_TEXT
    for ci, (txt, cw_c, cx_c) in enumerate(zip([comp, desc, corr], cws, cxs)):
        cell = rect(slide, cx_c, ry, cw_c, tr_h, fill=bg, line=NAVY, lw=Pt(0.6))
        tf2 = cell.text_frame; tf2.word_wrap = True
        p3  = tf2.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER if ci != 1 else PP_ALIGN.LEFT
        r3  = p3.add_run()
        r3.text = ("  " if ci == 1 else "") + txt
        r3.font.size = Pt(16 if not hdr else 17)
        r3.font.bold = hdr; r3.font.color.rgb = fg

# Population density note
note_y = tr_y0 + len(table_rows) * tr_h + 0.08
note_tf = textbox(slide, x2 + pad, note_y, C2W - 2 * pad, 0.45)
first_para(note_tf,
    "⚠  Population density (r ≈ −0.05) excluded: not an EPA formula input; negligible at city extent.",
    fs=14, color=MID_TEXT, bold=False)
first_para(note_tf, "", fs=1)  # dummy to avoid empty tf warning

y2 += calc_h + 0.4

# Key statistics
section_header(slide, x2, y2, C2W, "Key Statistics")
y2 += 0.46

sw   = (C2W - 0.2) / 2
sh   = 1.75
stats = [
    ("2,014",    "Census Tracts Analyzed"),
    ("18.4%",    "“Highly Walkable” or higher"),
    ("R² = 0.85", "Regression Model Fit"),
    ("p < 0.05", "Moran’s I: Significant Clustering"),
]
for i, (num, lbl) in enumerate(stats):
    r_i, c_i = divmod(i, 2)
    stat_box(slide, x2 + c_i * (sw + 0.2), y2 + r_i * (sh + 0.2), sw, sh, num, lbl)

y2 += 2 * sh + 0.2 * 2 + 0.4

# Corrplot — aspect ratio 1344:960 = 1.4:1
section_header(slide, x2, y2, C2W, "Correlation Matrix")
y2 += 0.46
corr_w = C2W
corr_h = corr_w / (1344 / 960)   # maintain aspect ratio → ≈ 9.6"
image(slide, "corrplot_debug.png", x2, y2, w=corr_w)
y2 += corr_h + 0.3

# Scatterplot matrix — aspect ratio 1600:1200 = 1.333:1
section_header(slide, x2, y2, C2W, "Scatterplot Matrix")
y2 += 0.46
scat_mat_w = C2W
image(slide, "poster_figures/fig_scatterplot_matrix.png", x2, y2, w=scat_mat_w)


# ── COLUMN 3: 2×2 map grid + scatter + conclusions ─────────────────────────
x3, y3 = C3X, BODY_Y

# Map aspect ratio: 2000×1600 → 1.25:1
MAP_RATIO   = 2000 / 1600   # 1.25
SCAT_RATIO  = 2000 / 1200   # 1.667

map_gap = 0.3
map_w   = (C3W - map_gap) / 2       # ≈ 11.05"
map_h   = map_w / MAP_RATIO          # ≈ 8.84"
cap_h   = 0.32

# ── Row 1: Greater Chicago + City Jenks ──
section_header(slide, x3, y3, C3W, "Walkability by Census Tract", fs=21)
y3 += 0.46

image(slide, "poster_figures/fig_metro_walk.png",  x3,            y3, w=map_w)
image(slide, "poster_figures/fig_city_jenks.png",  x3 + map_w + map_gap, y3, w=map_w)
y3 += map_h
caption(slide, x3,                    y3, map_w, "Greater Chicago Area (2021)")
caption(slide, x3 + map_w + map_gap,  y3, map_w, "City of Chicago (Adjusted Categories)")
y3 += cap_h + 0.25

# ── Row 2: Economic Opportunity + LISA ──
section_header(slide, x3, y3, C3W, "Economic Opportunity & Spatial Clusters (LISA)", fs=21)
y3 += 0.46

image(slide, "poster_figures/fig_econ_opportunity.png", x3,                    y3, w=map_w)
image(slide, "poster_figures/fig_lisa_city.png",        x3 + map_w + map_gap,  y3, w=map_w)
y3 += map_h
caption(slide, x3,                   y3, map_w, "Economic Opportunity (Lower Hardship)")
caption(slide, x3 + map_w + map_gap, y3, map_w, "LISA Clusters: City of Chicago (p < 0.05)")
y3 += cap_h + 0.25

# ── Scatter ──
section_header(slide, x3, y3, C3W, "Employment + Household Entropy vs. Walkability (City of Chicago)", fs=21)
y3 += 0.46

remaining = BODY_Y + BODY_H - y3 - 2.8   # leave room for conclusions
scat_h    = min(remaining, C3W / SCAT_RATIO)
scat_w    = scat_h * SCAT_RATIO
scat_x    = x3 + (C3W - scat_w) / 2
image(slide, "poster_figures/fig_scatter_entropy.png", scat_x, y3, w=scat_w)
y3 += scat_h + 0.3

# ── Conclusions ──
section_header(slide, x3, y3, C3W, "Conclusions", fs=21)
y3 += 0.46
tf = textbox(slide, x3, y3, C3W, BODY_Y + BODY_H - y3)
first_para(tf,
    "Walkability patterns strongly support the hypothesis, with the highest values "
    "concentrated in the Loop and decreasing outward. Moran's I confirms significant "
    "spatial clustering, and the strongest relationships are with EPA inputs like employment "
    "mix, entropy, and intersection density. The regression explains most variation "
    "(R² ≥ 0.85), while residential density shows almost no correlation, indicating "
    "that walkability is driven more by mixed-use development, connectivity, and transit "
    "access than by population density alone.\n"
    "These results align with planning ideas like the 15-minute city, where accessibility "
    "matters more than density. While the dataset strengthens ChiVes by adding "
    "built-environment insights, it does not capture factors like safety or sidewalk quality. "
    "Future work could refine the model, improve category definitions, and use machine "
    "learning to better predict walkability. Overall, improving walkability remains important "
    "for sustainability, equity, and urban quality of life.",
    fs=17, color=DARK_TEXT)


# ---------------------------------------------------------------------------
prs.save(str(BASE / "Undergraduate_Research_Symposium.pptx"))
print("Saved.")
